"""PPTX builder — turns a structured slide spec into a PowerPoint file.

Called by Dash's tool dispatch. The slide spec is a list of dicts, each with
a `layout` (title_slide, section, content, two_column, blank) and the text
content for that layout. Returns the path to the generated .pptx file.
"""

from __future__ import annotations

import os
import uuid
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT_DIR = Path(os.environ.get("DASH_OUTPUT_DIR", "/tmp/dash_output"))
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# Evergreen/Pokon brand-ish palette
_DARK = RGBColor(0x1B, 0x3A, 0x2D)
_ACCENT = RGBColor(0x2E, 0x7D, 0x32)
_LIGHT = RGBColor(0xE8, 0xF5, 0xE9)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_GRAY = RGBColor(0x61, 0x61, 0x61)


def _set_font(run, size_pt=14, bold=False, color=None):
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def _add_title_slide(prs: Presentation, spec: dict):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = _DARK

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = spec.get("title", "Forecast Review")
    _set_font(run, 36, bold=True, color=_WHITE)

    # Subtitle
    if spec.get("subtitle"):
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = spec["subtitle"]
        _set_font(run2, 18, color=_LIGHT)


def _add_section_slide(prs: Presentation, spec: dict):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = _ACCENT

    txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = spec.get("title", "")
    _set_font(run, 32, bold=True, color=_WHITE)


def _add_content_slide(prs: Presentation, spec: dict):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    # Heading
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(8.6), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = spec.get("heading", "")
    _set_font(run, 24, bold=True, color=_DARK)

    # Body bullets
    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.6), Inches(5))
    tf2 = body.text_frame
    tf2.word_wrap = True
    bullets = spec.get("bullets", [])
    if isinstance(bullets, str):
        bullets = [b.strip() for b in bullets.split("\n") if b.strip()]
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"• {bullet}" if not bullet.startswith("•") else bullet
        _set_font(run, 16, color=_GRAY)


def _add_two_column_slide(prs: Presentation, spec: dict):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    # Heading
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(8.6), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = spec.get("heading", "")
    _set_font(run, 24, bold=True, color=_DARK)

    for col_idx, col_key in enumerate(["left", "right"]):
        col = spec.get(col_key, {})
        x = Inches(0.7) if col_idx == 0 else Inches(5.2)
        box = slide.shapes.add_textbox(x, Inches(1.4), Inches(4.3), Inches(5))
        tf2 = box.text_frame
        tf2.word_wrap = True
        if col.get("title"):
            p = tf2.paragraphs[0]
            run = p.add_run()
            run.text = col["title"]
            _set_font(run, 18, bold=True, color=_ACCENT)
        bullets = col.get("bullets", [])
        if isinstance(bullets, str):
            bullets = [b.strip() for b in bullets.split("\n") if b.strip()]
        for bullet in bullets:
            p = tf2.add_paragraph()
            p.space_after = Pt(6)
            run = p.add_run()
            run.text = f"• {bullet}" if not bullet.startswith("•") else bullet
            _set_font(run, 14, color=_GRAY)


_BUILDERS = {
    "title_slide": _add_title_slide,
    "section": _add_section_slide,
    "content": _add_content_slide,
    "two_column": _add_two_column_slide,
}


def build_pptx(slides: list[dict[str, Any]], filename: str | None = None) -> dict:
    """Build a PPTX from a slide spec list. Returns {"path": ..., "filename": ...}."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for spec in slides:
        layout = spec.get("layout", "content")
        builder = _BUILDERS.get(layout, _add_content_slide)
        builder(prs, spec)

    fname = filename or f"dash_{uuid.uuid4().hex[:8]}.pptx"
    if not fname.endswith(".pptx"):
        fname += ".pptx"
    out_path = OUTPUT_DIR / fname
    prs.save(str(out_path))
    return {"path": str(out_path), "filename": fname, "slides": len(slides)}
